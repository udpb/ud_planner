"""
Build the Underdogs 1N2D Hackathon lecture deck (44 slides).
Output: deliverables/lecture-deck.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Brand tokens ----------
ORANGE = RGBColor(0xF0, 0x55, 0x19)
ORANGE_SOFT = RGBColor(0xF4, 0x80, 0x53)
ORANGE_PALE = RGBColor(0xFF, 0xF5, 0xF0)
CYAN = RGBColor(0x06, 0xA9, 0xD0)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK = RGBColor(0x37, 0x39, 0x38)
GRAY_BODY = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
BORDER = RGBColor(0xD8, 0xD4, 0xD7)
CARD_BG = RGBColor(0xF5, 0xF0, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Inter"
TITLE_SIZE = Pt(24)
SECTION_SIZE = Pt(16)
BODY_SIZE = Pt(11)
SMALL_SIZE = Pt(9)
TINY_SIZE = Pt(8)
MICRO_SIZE = Pt(7)

SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.625)
MARGIN = Inches(0.5)
CONTENT_W = SLIDE_W - (MARGIN * 2)

COPYRIGHT = "(c) 2026 UD Impact Co., Ltd. (Underdogs)"

# ---------- Low-level helpers ----------

def _no_fill(shape):
    shape.fill.background()

def _solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def _no_line(shape):
    shape.line.fill.background()

def _line(shape, rgb, width_pt=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill=None, border=None, border_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        _no_fill(shp)
    else:
        _solid(shp, fill)
    if border is None:
        _no_line(shp)
    else:
        _line(shp, border, border_w)
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=None,
    bold=False,
    italic=False,
    color=BLACK,
    font=FONT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        runs = [text]
    else:
        runs = text

    for i, item in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, str):
            seg = item
            run_opts = {}
        else:
            seg = item.get("text", "")
            run_opts = item
        p.alignment = run_opts.get("align", align)
        p.line_spacing = run_opts.get("line_spacing", line_spacing)

        r = p.add_run()
        r.text = seg
        rf = r.font
        rf.name = run_opts.get("font", font)
        rf.size = run_opts.get("size", size if size else BODY_SIZE)
        rf.bold = run_opts.get("bold", bold)
        rf.italic = run_opts.get("italic", italic)
        rf.color.rgb = run_opts.get("color", color)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=None, color=BLACK, spacing=1.25, bullet_color=ORANGE):
    """Small custom bullet list using bullet character."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing

        r1 = p.add_run()
        r1.text = "\u25A0  "  # small square
        r1.font.name = FONT
        r1.font.size = size or BODY_SIZE
        r1.font.bold = True
        r1.font.color.rgb = bullet_color

        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = size or BODY_SIZE
        r2.font.color.rgb = color
    return tb


# ---------- Slide scaffolding ----------

def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Paint white background explicitly
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    return slide


def add_footer(slide, section_name, page_num):
    # Bottom copyright (right, 7pt gray)
    add_text(
        slide,
        Inches(0.5),
        Inches(5.32),
        Inches(9.0),
        Inches(0.2),
        COPYRIGHT,
        size=MICRO_SIZE,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )
    # Section label (left, 8pt gray)
    if section_name:
        add_text(
            slide,
            Inches(0.5),
            Inches(5.15),
            Inches(5.0),
            Inches(0.2),
            section_name,
            size=TINY_SIZE,
            color=GRAY,
        )
    # Page number (right, 8pt orange bold)
    add_text(
        slide,
        Inches(8.5),
        Inches(5.15),
        Inches(1.0),
        Inches(0.2),
        str(page_num).zfill(2),
        size=TINY_SIZE,
        bold=True,
        color=ORANGE,
        align=PP_ALIGN.RIGHT,
    )


def add_title_bar(slide, title, eyebrow=None):
    """Standard slide title area with optional eyebrow label."""
    y = Inches(0.45)
    if eyebrow:
        add_text(
            slide,
            MARGIN,
            y,
            CONTENT_W,
            Inches(0.22),
            eyebrow.upper(),
            size=SMALL_SIZE,
            bold=True,
            color=ORANGE,
        )
        y = Inches(0.72)
    add_text(
        slide,
        MARGIN,
        y,
        CONTENT_W,
        Inches(0.5),
        title,
        size=TITLE_SIZE,
        bold=True,
        color=BLACK,
    )
    # Thin orange underline accent (subtle, short)
    add_rect(
        slide,
        MARGIN,
        Inches(1.25),
        Inches(0.5),
        Inches(0.04),
        fill=ORANGE,
    )


# ---------- Page templates ----------

def section_divider(prs, part_label, headline, blurb, page_num, section_name):
    slide = new_slide(prs)
    # Black background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BLACK)
    # Left orange accent bar
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(0.08), Inches(2.8), fill=ORANGE)
    # Eyebrow label
    add_text(
        slide,
        Inches(0.85),
        Inches(1.3),
        Inches(8.0),
        Inches(0.3),
        part_label.upper(),
        size=Pt(10),
        bold=True,
        color=ORANGE,
    )
    # Headline
    add_text(
        slide,
        Inches(0.85),
        Inches(1.65),
        Inches(8.3),
        Inches(1.5),
        headline,
        size=Pt(34),
        bold=True,
        color=WHITE,
        line_spacing=1.1,
    )
    # Blurb
    add_text(
        slide,
        Inches(0.85),
        Inches(3.35),
        Inches(8.3),
        Inches(1.0),
        blurb,
        size=Pt(12),
        color=RGBColor(0xD8, 0xD4, 0xD7),
        line_spacing=1.35,
    )
    # Page num
    add_text(
        slide,
        Inches(8.5),
        Inches(5.15),
        Inches(1.0),
        Inches(0.2),
        str(page_num).zfill(2),
        size=TINY_SIZE,
        bold=True,
        color=ORANGE,
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(5.15),
        Inches(5.0),
        Inches(0.2),
        section_name,
        size=TINY_SIZE,
        color=RGBColor(0x99, 0x99, 0x99),
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(5.32),
        Inches(9.0),
        Inches(0.2),
        COPYRIGHT,
        size=MICRO_SIZE,
        color=RGBColor(0x99, 0x99, 0x99),
        align=PP_ALIGN.RIGHT,
    )
    return slide


def three_column_cards(slide, y, cards, *, card_h=Inches(3.0), gap=Inches(0.15)):
    """Render 3 (or more) cards with orange left bar + pale bg."""
    n = len(cards)
    total_gap = gap * (n - 1)
    card_w = (CONTENT_W - total_gap) / n
    for i, c in enumerate(cards):
        cx = MARGIN + (card_w + gap) * i
        # Background card
        add_rect(slide, cx, y, card_w, card_h, fill=CARD_BG, border=BORDER, border_w=0.5)
        # Orange left bar
        add_rect(slide, cx, y, Inches(0.06), card_h, fill=ORANGE)
        # Step/tag label
        if c.get("tag"):
            add_text(
                slide,
                cx + Inches(0.22),
                y + Inches(0.18),
                card_w - Inches(0.35),
                Inches(0.22),
                c["tag"].upper(),
                size=SMALL_SIZE,
                bold=True,
                color=ORANGE,
            )
        # Title
        add_text(
            slide,
            cx + Inches(0.22),
            y + Inches(0.45),
            card_w - Inches(0.35),
            Inches(0.5),
            c["title"],
            size=Pt(13),
            bold=True,
            color=BLACK,
            line_spacing=1.15,
        )
        # Body
        body_y = y + Inches(1.0)
        body_h = card_h - Inches(1.1)
        body = c.get("body")
        if isinstance(body, list):
            add_bullets(
                slide,
                cx + Inches(0.22),
                body_y,
                card_w - Inches(0.35),
                body_h,
                body,
                size=Pt(10),
                color=GRAY_BODY,
                spacing=1.2,
            )
        elif body:
            add_text(
                slide,
                cx + Inches(0.22),
                body_y,
                card_w - Inches(0.35),
                body_h,
                body,
                size=Pt(10),
                color=GRAY_BODY,
                line_spacing=1.3,
            )


def stat_callout(slide, x, y, w, h, number, label):
    add_rect(slide, x, y, w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
    add_rect(slide, x, y, Inches(0.06), h, fill=ORANGE)
    add_text(
        slide,
        x + Inches(0.2),
        y + Inches(0.15),
        w - Inches(0.3),
        Inches(0.7),
        number,
        size=Pt(26),
        bold=True,
        color=ORANGE,
    )
    add_text(
        slide,
        x + Inches(0.2),
        y + Inches(0.85),
        w - Inches(0.3),
        h - Inches(0.95),
        label,
        size=Pt(9),
        color=GRAY_BODY,
        line_spacing=1.25,
    )


def two_col_text(slide, y, h, left_title, left_body, right_title, right_body):
    col_w = (CONTENT_W - Inches(0.3)) / 2
    # Left
    add_rect(slide, MARGIN, y, col_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
    add_rect(slide, MARGIN, y, Inches(0.06), h, fill=ORANGE)
    add_text(slide, MARGIN + Inches(0.22), y + Inches(0.2), col_w - Inches(0.35), Inches(0.35),
             left_title, size=Pt(13), bold=True, color=BLACK)
    if isinstance(left_body, list):
        add_bullets(slide, MARGIN + Inches(0.22), y + Inches(0.62), col_w - Inches(0.35), h - Inches(0.75),
                    left_body, size=Pt(10), color=GRAY_BODY, spacing=1.25)
    else:
        add_text(slide, MARGIN + Inches(0.22), y + Inches(0.62), col_w - Inches(0.35), h - Inches(0.75),
                 left_body, size=Pt(10), color=GRAY_BODY, line_spacing=1.35)
    # Right
    rx = MARGIN + col_w + Inches(0.3)
    add_rect(slide, rx, y, col_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
    add_rect(slide, rx, y, Inches(0.06), h, fill=ORANGE)
    add_text(slide, rx + Inches(0.22), y + Inches(0.2), col_w - Inches(0.35), Inches(0.35),
             right_title, size=Pt(13), bold=True, color=BLACK)
    if isinstance(right_body, list):
        add_bullets(slide, rx + Inches(0.22), y + Inches(0.62), col_w - Inches(0.35), h - Inches(0.75),
                    right_body, size=Pt(10), color=GRAY_BODY, spacing=1.25)
    else:
        add_text(slide, rx + Inches(0.22), y + Inches(0.62), col_w - Inches(0.35), h - Inches(0.75),
                 right_body, size=Pt(10), color=GRAY_BODY, line_spacing=1.35)


# ---------- Build ----------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # =============================================================
    # PART A — Opening (slides 1-3)
    # =============================================================

    # Slide 1 — Cover
    s = new_slide(prs)
    # Full pale bg
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=ORANGE_PALE)
    # Right-side orange block (motif: solid slab)
    add_rect(s, Inches(6.8), 0, Inches(3.2), SLIDE_H, fill=ORANGE)
    # Vertical thin black accent
    add_rect(s, Inches(6.76), 0, Inches(0.04), SLIDE_H, fill=BLACK)
    # Eyebrow
    add_text(s, Inches(0.6), Inches(0.85), Inches(6.0), Inches(0.3),
             "UNDERDOGS EDUCATION PLANNING LECTURE", size=Pt(10), bold=True, color=ORANGE)
    # Title
    add_text(s, Inches(0.6), Inches(1.25), Inches(6.0), Inches(2.2),
             "Designing\n1-Night-2-Day\nHackathons", size=Pt(36), bold=True, color=BLACK,
             line_spacing=1.05)
    # Subtitle
    add_text(s, Inches(0.6), Inches(3.55), Inches(6.0), Inches(0.9),
             "Underdogs Education Planning for\nIndonesia \u00B7 Korea \u00B7 Japan \u00B7 Taiwan \u00B7 India",
             size=Pt(13), color=GRAY_BODY, line_spacing=1.35)
    # Right block: version / audience
    add_text(s, Inches(7.05), Inches(1.25), Inches(2.7), Inches(0.4),
             "FOR", size=Pt(9), bold=True, color=WHITE)
    add_text(s, Inches(7.05), Inches(1.5), Inches(2.7), Inches(1.5),
             "External Coaches\n+ Planning Team\n+ 5-Country Partners", size=Pt(13), bold=True,
             color=WHITE, line_spacing=1.3)
    add_text(s, Inches(7.05), Inches(3.8), Inches(2.7), Inches(0.3),
             "VERSION", size=Pt(9), bold=True, color=WHITE)
    add_text(s, Inches(7.05), Inches(4.05), Inches(2.7), Inches(0.3),
             "v1.0  /  April 2026", size=Pt(12), bold=True, color=WHITE)
    # Bottom copyright only (no footer number on cover)
    add_text(s, Inches(0.6), Inches(5.3), Inches(6.0), Inches(0.2),
             COPYRIGHT, size=MICRO_SIZE, color=GRAY)

    # Slide 2 — Why This Matters
    s = new_slide(prs)
    add_title_bar(s, "Short-format hackathons are the fastest way to install founder mindset.",
                  eyebrow="Why this matters \u00B7 Workshop goal")
    add_text(s, MARGIN, Inches(1.45), CONTENT_W, Inches(0.5),
             "By the end of this lecture you can design a 1-night-2-day hackathon that reflects Underdogs' philosophy, localizes to country context, and connects pre-event LMS with post-event batch.",
             size=Pt(12), color=GRAY_BODY, line_spacing=1.35)
    cards = [
        {"tag": "The problem", "title": "Long programs lose momentum",
         "body": ["3-4 month batches suffer drop-off",
                  "Theory without action kills learning",
                  "Partners want fast proof of concept"]},
        {"tag": "Our wager", "title": "36 hours, one competency",
         "body": ["One 5D competency per hackathon",
                  "Every 2 hours = visible deliverable",
                  "Pre-LMS + 1N2D + post-batch = funnel"]},
        {"tag": "Your outcome", "title": "A run-sheet you can ship",
         "body": ["8-block template, localized",
                  "Coach cue cards per axis",
                  "Demo day scorecard ready"]},
    ]
    three_column_cards(s, Inches(2.45), cards, card_h=Inches(2.55))
    add_footer(s, "Part A \u00B7 Opening", 2)

    # Slide 3 — Agenda
    s = new_slide(prs)
    add_title_bar(s, "Six parts, forty-four slides, five country tracks.",
                  eyebrow="Agenda")
    items = [
        ("A", "Opening", "Why 1N2D, goal, agenda"),
        ("B", "Underdogs Planning Lens", "Philosophy \u00B7 Impact Value Chain \u00B7 9 methodologies \u00B7 5D"),
        ("C", "Why 1N2D Works", "Funnel \u00B7 8-block template \u00B7 coaching cadence \u00B7 rubric"),
        ("D", "5 Country Tracks", "Indonesia \u00B7 KR-Inbound \u00B7 KR-Outbound \u00B7 JP+TW \u00B7 India"),
        ("E", "Facilitator Playbook", "Pre-event \u00B7 Day 1 \u00B7 Day 2 \u00B7 cue cards \u00B7 post-event"),
        ("F", "Close", "Takeaways \u00B7 resources \u00B7 Q&A"),
    ]
    row_h = Inches(0.55)
    base_y = Inches(1.55)
    for i, (letter, name, desc) in enumerate(items):
        ry = base_y + row_h * i
        # Letter badge
        add_rect(s, MARGIN, ry, Inches(0.55), Inches(0.45), fill=ORANGE)
        add_text(s, MARGIN, ry, Inches(0.55), Inches(0.45),
                 letter, size=Pt(18), bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_text(s, MARGIN + Inches(0.75), ry + Inches(0.02), Inches(3.2), Inches(0.3),
                 name, size=Pt(13), bold=True, color=BLACK)
        # Description
        add_text(s, MARGIN + Inches(0.75), ry + Inches(0.28), Inches(7.5), Inches(0.22),
                 desc, size=Pt(10), color=GRAY_BODY)
        # Divider line
        if i < len(items) - 1:
            add_rect(s, MARGIN, ry + row_h - Inches(0.02), CONTENT_W, Inches(0.008), fill=BORDER)
    add_footer(s, "Part A \u00B7 Opening", 3)

    # =============================================================
    # PART B — Planning Lens (slides 4-12)
    # =============================================================

    # Slide 4 — Part B divider
    section_divider(prs,
        "Part B",
        "The Underdogs planning lens.",
        "Six principles, one value chain, nine methodologies, five competencies, four tiers of support. Every hackathon we design carries this DNA.",
        4, "Part B \u00B7 Planning Lens")

    # Slide 5 — Six Principles
    s = new_slide(prs)
    add_title_bar(s, "Six principles guide every Underdogs program design.",
                  eyebrow="Planning philosophy")
    principles = [
        ("01", "Data flows top-down", "Each step inherits prior artifacts via PipelineContext."),
        ("02", "Assets rise automatically", "IMPACT modules, coaches, SROI proxies auto-surface."),
        ("03", "AI operates in context", "Never from scratch\u2014always on accumulated signal."),
        ("04", "New PMs know the why", "Guides, references, warnings embedded per step."),
        ("05", "Impact-First on curriculum", "Activities auto-extract; AI generates Outcome + Impact."),
        ("06", "Action Week is non-negotiable", "3 theory sessions max before a practice week."),
    ]
    col_w = (CONTENT_W - Inches(0.3)) / 3
    row_h = Inches(1.65)
    for i, (num, title, desc) in enumerate(principles):
        r = i // 3
        c = i % 3
        x = MARGIN + (col_w + Inches(0.15)) * c
        y = Inches(1.5) + (row_h + Inches(0.15)) * r
        add_rect(s, x, y, col_w, row_h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), row_h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.3), Inches(0.3),
                 num, size=Pt(11), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.4), col_w - Inches(0.3), Inches(0.5),
                 title, size=Pt(12), bold=True, color=BLACK, line_spacing=1.15)
        add_text(s, x + Inches(0.2), y + Inches(0.95), col_w - Inches(0.3), row_h - Inches(1.05),
                 desc, size=Pt(9), color=GRAY_BODY, line_spacing=1.3)
    add_footer(s, "Part B \u00B7 Planning Lens", 5)

    # Slide 6 — Impact Value Chain
    s = new_slide(prs)
    add_title_bar(s, "Input \u2192 Activity \u2192 Output \u2192 Outcome \u2192 Impact. Five layers, one causal line.",
                  eyebrow="Impact value chain")
    stages = [
        ("Input", "Resources", "Budget \u00B7 coaches \u00B7 LMS \u00B7 venue"),
        ("Activity", "What we do", "1N2D blocks \u00B7 1:1 coaching \u00B7 demo"),
        ("Output", "What is produced", "MVPs \u00B7 pitch decks \u00B7 live URLs"),
        ("Outcome", "What changes", "5D competency up \u00B7 team formed"),
        ("Impact", "Why it matters", "Problem solved \u00B7 market moved"),
    ]
    seg_w = (CONTENT_W - Inches(0.4)) / 5
    y = Inches(1.8)
    h = Inches(2.3)
    for i, (stage, tag, body) in enumerate(stages):
        x = MARGIN + (seg_w + Inches(0.1)) * i
        add_rect(s, x, y, seg_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, seg_w, Inches(0.42), fill=ORANGE)
        add_text(s, x, y, seg_w, Inches(0.42), stage.upper(),
                 size=Pt(12), bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), y + Inches(0.58), seg_w - Inches(0.25), Inches(0.3),
                 tag, size=Pt(10), bold=True, color=BLACK)
        add_text(s, x + Inches(0.15), y + Inches(0.9), seg_w - Inches(0.25), h - Inches(1.0),
                 body, size=Pt(9), color=GRAY_BODY, line_spacing=1.3)
        # Arrow between segments
        if i < len(stages) - 1:
            ax = x + seg_w + Inches(0.005)
            add_text(s, ax, y + Inches(0.9), Inches(0.1), Inches(0.4),
                     "\u25B8", size=Pt(14), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(4.35), CONTENT_W, Inches(0.4),
             "5D mapping: Desire + Direction feed Outcome/Impact (why). Discipline + Design feed Activity/Output (how). Durability feeds long-term Impact.",
             size=Pt(10), italic=True, color=GRAY_BODY, line_spacing=1.35)
    add_footer(s, "Part B \u00B7 Planning Lens", 6)

    # Slide 7 — So What? Test
    s = new_slide(prs)
    add_title_bar(s, 'Every idea must survive three "So what?" rounds before pitching.',
                  eyebrow='The "so what?" test')
    cards = [
        {"tag": "Round 1", "title": "So what... problem?",
         "body": ["Who exactly suffers?",
                  "How often, how much?",
                  "Why now, not 5 years ago?"]},
        {"tag": "Round 2", "title": "So what... solution?",
         "body": ["What's the non-obvious wedge?",
                  "Why you vs. incumbents?",
                  "Which assumption is riskiest?"]},
        {"tag": "Round 3", "title": "So what... impact?",
         "body": ["Market size \u00D7 social ripple",
                  "Tied to Outcome layer",
                  "Measurable within 12 months"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.6))
    add_text(s, MARGIN, Inches(4.3), CONTENT_W, Inches(0.4),
             "If a team cannot answer all three, send them back to the problem deep-dive. Coach time is too expensive to waste on under-specified ideas.",
             size=Pt(10), italic=True, color=GRAY_BODY, line_spacing=1.35)
    add_footer(s, "Part B \u00B7 Planning Lens", 7)

    # Slide 8 — 9 Methodologies
    s = new_slide(prs)
    add_title_bar(s, "Nine methodologies. You pick, you don't blend.",
                  eyebrow="Methodology is not one-size")
    methods = [
        ("IMPACT", "18 modules", "Flagship. Impact-first curriculum."),
        ("Local", "Region-fit", "City / province / community anchors."),
        ("Glocal", "Local to global", "Local problem, global scaling path."),
        ("Competition", "Contest mode", "Hackathons, demo-day pitches."),
        ("Matching", "Team assembly", "Cross-discipline cofounder matching."),
        ("Re-founding", "2nd act", "Pivot / exit-driven founders."),
        ("Global Entry", "Market entry", "Cross-border expansion tracks."),
        ("Small Business", "Sosang-gongin", "B2C / local economy operators."),
        ("Custom", "Partner-fit", "Tailored per partner brief."),
    ]
    col_w = (CONTENT_W - Inches(0.3)) / 3
    row_h = Inches(0.95)
    for i, (name, tag, desc) in enumerate(methods):
        r = i // 3
        c = i % 3
        x = MARGIN + (col_w + Inches(0.15)) * c
        y = Inches(1.5) + (row_h + Inches(0.1)) * r
        add_rect(s, x, y, col_w, row_h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), row_h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
                 name, size=Pt(12), bold=True, color=BLACK)
        add_text(s, x + Inches(0.2), y + Inches(0.36), col_w - Inches(0.3), Inches(0.22),
                 tag.upper(), size=Pt(8), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.58), col_w - Inches(0.3), Inches(0.35),
                 desc, size=Pt(9), color=GRAY_BODY, line_spacing=1.25)
    add_footer(s, "Part B \u00B7 Planning Lens", 8)

    # Slide 9 — 5D Competencies
    s = new_slide(prs)
    add_title_bar(s, "5D competencies \u2014 we measure people, not just products.",
                  eyebrow="Five dimensions (ACTT extension)")
    rows = [
        ("D1", "Desire", "Why do you want this?", "Clear personal stake / story"),
        ("D2", "Direction", "Where are you going?", "Specific goal + near-term milestone"),
        ("D3", "Discipline", "How do you execute?", "Habit + cadence + self-monitoring"),
        ("D4", "Design", "How do you prototype?", "Tests assumption with MVP"),
        ("D5", "Durability", "How do you sustain?", "Resilience + learning loop + team"),
    ]
    # Header row
    y = Inches(1.5)
    row_h = Inches(0.55)
    headers = [("Code", Inches(0.6)), ("Name", Inches(1.15)), ("Core question", Inches(3.4)), ("Observable signal", Inches(3.9))]
    x = MARGIN
    add_rect(s, MARGIN, y, CONTENT_W, Inches(0.38), fill=BLACK)
    for label, w in headers:
        add_text(s, x + Inches(0.15), y + Inches(0.05), w, Inches(0.3),
                 label, size=Pt(10), bold=True, color=WHITE)
        x += w
    y += Inches(0.38)
    for i, (code, name, q, sig) in enumerate(rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_rect(s, MARGIN, y, CONTENT_W, row_h, fill=bg, border=BORDER, border_w=0.3)
        x = MARGIN
        add_text(s, x + Inches(0.15), y + Inches(0.15), Inches(0.6), row_h - Inches(0.2),
                 code, size=Pt(11), bold=True, color=ORANGE)
        x += Inches(0.6)
        add_text(s, x + Inches(0.15), y + Inches(0.15), Inches(1.15), row_h - Inches(0.2),
                 name, size=Pt(11), bold=True, color=BLACK)
        x += Inches(1.15)
        add_text(s, x + Inches(0.15), y + Inches(0.15), Inches(3.4), row_h - Inches(0.2),
                 q, size=Pt(10), color=GRAY_BODY)
        x += Inches(3.4)
        add_text(s, x + Inches(0.15), y + Inches(0.15), Inches(3.9), row_h - Inches(0.2),
                 sig, size=Pt(10), color=GRAY_BODY)
        y += row_h
    add_text(s, MARGIN, y + Inches(0.15), CONTENT_W, Inches(0.3),
             "1N2D rule: install ONE D per hackathon. You cannot teach all five in 36 hours.",
             size=Pt(10), italic=True, color=ORANGE, bold=True)
    add_footer(s, "Part B \u00B7 Planning Lens", 9)

    # Slide 10 — 4-Tier Support
    s = new_slide(prs)
    add_title_bar(s, "Four tiers of support \u2014 apply only when the stage calls for it.",
                  eyebrow="4-tier support system")
    cards = [
        {"tag": "Tier 1", "title": "Cohort education",
         "body": ["Async LMS + live workshops",
                  "Applies to ALL participants",
                  "1N2D: pre-event 2-3 weeks"]},
        {"tag": "Tier 2", "title": "1:1 coaching",
         "body": ["3-4 sessions per team",
                  "20-30 min per session",
                  "Applies to selected teams"]},
        {"tag": "Tier 3", "title": "Batch incubation",
         "body": ["3-4 month follow-up",
                  "Top 10-20% of hackathon teams",
                  "Milestones + graduation criteria"]},
        {"tag": "Tier 4", "title": "Alliance / fund",
         "body": ["VC + partner network",
                  "Flagship brand treatment",
                  "For graduates of Tier 3"]},
    ]
    col_w = (CONTENT_W - Inches(0.3)) / 4
    y = Inches(1.55)
    h = Inches(3.0)
    for i, c in enumerate(cards):
        x = MARGIN + (col_w + Inches(0.1)) * i
        add_rect(s, x, y, col_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.18), col_w - Inches(0.3), Inches(0.22),
                 c["tag"].upper(), size=Pt(9), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.45), col_w - Inches(0.3), Inches(0.45),
                 c["title"], size=Pt(12), bold=True, color=BLACK, line_spacing=1.15)
        add_bullets(s, x + Inches(0.2), y + Inches(0.95), col_w - Inches(0.3), h - Inches(1.05),
                    c["body"], size=Pt(9), color=GRAY_BODY, spacing=1.25)
    add_footer(s, "Part B \u00B7 Planning Lens", 10)

    # Slide 11 — Action Week
    s = new_slide(prs)
    add_title_bar(s, "Three theory sessions in a row trigger an Action Week. No exceptions.",
                  eyebrow="Execution-guaranteed: action week")
    # Left: the rule
    add_rect(s, MARGIN, Inches(1.55), Inches(4.6), Inches(3.0), fill=BLACK)
    add_rect(s, MARGIN, Inches(1.55), Inches(0.08), Inches(3.0), fill=ORANGE)
    add_text(s, MARGIN + Inches(0.3), Inches(1.7), Inches(4.2), Inches(0.35),
             "THE RULE", size=Pt(10), bold=True, color=ORANGE)
    add_text(s, MARGIN + Inches(0.3), Inches(2.05), Inches(4.2), Inches(1.8),
             "If curriculum stacks 3 consecutive theory sessions, insert Action Week \u2014 a hands-on block where teams apply, not absorb.",
             size=Pt(14), bold=True, color=WHITE, line_spacing=1.3)
    add_text(s, MARGIN + Inches(0.3), Inches(3.8), Inches(4.2), Inches(0.6),
             "The system warns the PM automatically when the trigger fires.",
             size=Pt(10), color=RGBColor(0xD8, 0xD4, 0xD7), line_spacing=1.35)
    # Right: 1N2D translation
    rx = MARGIN + Inches(4.8)
    rw = CONTENT_W - Inches(4.8)
    add_rect(s, rx, Inches(1.55), rw, Inches(3.0), fill=CARD_BG, border=BORDER, border_w=0.5)
    add_rect(s, rx, Inches(1.55), Inches(0.06), Inches(3.0), fill=ORANGE)
    add_text(s, rx + Inches(0.2), Inches(1.7), rw - Inches(0.3), Inches(0.3),
             "IN A 1N2D", size=Pt(10), bold=True, color=ORANGE)
    add_bullets(s, rx + Inches(0.2), Inches(2.05), rw - Inches(0.3), Inches(2.4),
                ["Max 45 min lecture per 3-hour block",
                 "Each block must close with a visible deliverable",
                 "Coaching time \u2265 lecture time",
                 "If concept-heavy, move to pre-LMS",
                 "On-site time is for building, not listening"],
                size=Pt(11), color=GRAY_BODY, spacing=1.35)
    add_footer(s, "Part B \u00B7 Planning Lens", 11)

    # Slide 12 — Quantitative Saturation + Section V
    s = new_slide(prs)
    add_title_bar(s, "Numbers everywhere. No claim without a figure.",
                  eyebrow="Quantitative saturation + section v bonus")
    # Stat row
    stats = [
        ("520+", "Global partners"),
        ("800", "Coach pool"),
        ("18", "IMPACT modules"),
        ("9", "Methodologies"),
        ("5", "5D competencies"),
        ("4", "Support tiers"),
    ]
    y = Inches(1.55)
    sw = (CONTENT_W - Inches(0.5)) / 6
    for i, (num, label) in enumerate(stats):
        x = MARGIN + (sw + Inches(0.1)) * i
        stat_callout(s, x, y, sw, Inches(1.3), num, label)
    # Section V note
    add_rect(s, MARGIN, Inches(3.15), CONTENT_W, Inches(1.55), fill=BLACK)
    add_rect(s, MARGIN, Inches(3.15), Inches(0.08), Inches(1.55), fill=ORANGE)
    add_text(s, MARGIN + Inches(0.3), Inches(3.3), CONTENT_W - Inches(0.5), Inches(0.35),
             "SECTION V BONUS", size=Pt(10), bold=True, color=ORANGE)
    add_text(s, MARGIN + Inches(0.3), Inches(3.6), CONTENT_W - Inches(0.5), Inches(1.0),
             "Top proposals reserve a Section V \u2014 beyond-the-brief moves: new partners, new tracks, new validation instruments. Judges remember the team that exceeded scope.",
             size=Pt(12), color=WHITE, line_spacing=1.35)
    add_footer(s, "Part B \u00B7 Planning Lens", 12)

    # =============================================================
    # PART C — Why 1N2D Works (slides 13-19)
    # =============================================================

    # Slide 13 — Part C divider
    section_divider(prs,
        "Part C",
        "Why 1-night-2-day works.",
        "Short-format forces focus. The funnel makes it rigorous. The 8-block template makes it repeatable. The rubric makes it fair.",
        13, "Part C \u00B7 Why 1N2D")

    # Slide 14 — Short vs Long
    s = new_slide(prs)
    add_title_bar(s, "Short-format trades depth for momentum \u2014 and momentum wins at Stage 1.",
                  eyebrow="Short-format vs long-format")
    two_col_text(s, Inches(1.55), Inches(3.3),
                 "Long-format (3-4 months)",
                 ["Deep validation, real customers",
                  "Drop-off 40-60% common",
                  "Commitment bar high",
                  "Requires proven teams",
                  "Expensive per-participant"],
                 "Short-format (1N2D \u00B7 36h)",
                 ["Forced completion, visible output",
                  "Drop-off near zero",
                  "Low commitment to enter",
                  "Reveals 5D signals quickly",
                  "Cheap to run at scale"])
    add_text(s, MARGIN, Inches(5.0), CONTENT_W, Inches(0.3),
             "Underdogs uses both. 1N2D is the screening + activation layer; the batch program does the depth work.",
             size=Pt(10), italic=True, color=GRAY_BODY)
    add_footer(s, "Part C \u00B7 Why 1N2D", 14)

    # Slide 15 — Funnel
    s = new_slide(prs)
    add_title_bar(s, "Pre-LMS \u2192 1N2D \u2192 Follow-up batch. The hackathon is the middle, not the whole.",
                  eyebrow="The funnel")
    stages = [
        ("Pre-event LMS", "2-3 weeks async",
         ["5-8 video modules",
          "3-5 readings",
          "Pre-arrival exercise",
          "Team formation hint"]),
        ("1N2D hackathon", "36 hours on-site",
         ["8 blocks Day 1 + 2",
          "3-4 coach 1:1s per team",
          "Demo day + scorecard",
          "Awards ceremony"]),
        ("Follow-up batch", "2-8 weeks \u2192 3-4 months",
         ["1-pager report every team",
          "Top 10-20% \u2192 Stage 2 batch",
          "Alumni quarterly check-in",
          "Alliance / fund for Stage 3"]),
    ]
    y = Inches(1.55)
    seg_w = (CONTENT_W - Inches(0.4)) / 3
    h = Inches(3.1)
    for i, (name, dur, body) in enumerate(stages):
        x = MARGIN + (seg_w + Inches(0.2)) * i
        add_rect(s, x, y, seg_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, seg_w, Inches(0.5), fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.08), seg_w - Inches(0.3), Inches(0.3),
                 name, size=Pt(13), bold=True, color=WHITE)
        add_text(s, x + Inches(0.2), y + Inches(0.75), seg_w - Inches(0.3), Inches(0.25),
                 dur.upper(), size=Pt(9), bold=True, color=ORANGE)
        add_bullets(s, x + Inches(0.2), y + Inches(1.1), seg_w - Inches(0.3), h - Inches(1.2),
                    body, size=Pt(10), color=GRAY_BODY, spacing=1.3)
        if i < 2:
            ax = x + seg_w + Inches(0.02)
            add_text(s, ax, y + Inches(1.35), Inches(0.18), Inches(0.5),
                     "\u25B8", size=Pt(18), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_footer(s, "Part C \u00B7 Why 1N2D", 15)

    # Slide 16 — Engagement techniques
    s = new_slide(prs)
    add_title_bar(s, "Seven techniques keep energy high across 36 hours.",
                  eyebrow="Engagement techniques for 24+ hours")
    techniques = [
        ("01", "Milestone compression", "Every 2-hour block has a visible deliverable."),
        ("02", "1:1 coaching density", "3-4 sessions per team, 20-30 min each."),
        ("03", "Peer accountability", "Public daily standup \u2014 2 min per team."),
        ("04", "Immersive environment", "Dinner together, optional overnight on-site."),
        ("05", "Physical artifact", "Every team produces tangible output (poster / URL)."),
        ("06", "Public demo", "External judges, not just internal coaches."),
        ("07", "Celebration ritual", "Award ceremony, group photo, post-event social."),
    ]
    col_w = (CONTENT_W - Inches(0.3)) / 3
    row_h = Inches(1.1)
    for i, (num, title, desc) in enumerate(techniques):
        r = i // 3
        c = i % 3
        x = MARGIN + (col_w + Inches(0.15)) * c
        y = Inches(1.55) + (row_h + Inches(0.15)) * r
        add_rect(s, x, y, col_w, row_h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), row_h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.12), col_w - Inches(0.3), Inches(0.22),
                 num, size=Pt(10), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.36), col_w - Inches(0.3), Inches(0.3),
                 title, size=Pt(12), bold=True, color=BLACK)
        add_text(s, x + Inches(0.2), y + Inches(0.66), col_w - Inches(0.3), row_h - Inches(0.75),
                 desc, size=Pt(9), color=GRAY_BODY, line_spacing=1.3)
    add_footer(s, "Part C \u00B7 Why 1N2D", 16)

    # Slide 17 — 8-Block Template
    s = new_slide(prs)
    add_title_bar(s, "Eight blocks \u2014 five on Day 1, three on Day 2. Every block closes with an artifact.",
                  eyebrow="The 8-block 1N2D template")
    blocks = [
        ("B1", "D1 \u00B7 09-10", "Opening + Icebreaker", "Team norms + intros"),
        ("B2", "D1 \u00B7 10-12", "Problem deep-dive", "Problem statement v1"),
        ("B3", "D1 \u00B7 13-15", "Team + Ideation", "Team locked + idea v1"),
        ("B4", "D1 \u00B7 15-18", "MVP sprint #1 + 1:1", "First prototype"),
        ("B5", "D1 \u00B7 19-22", "Evening sprint + dinner", "Rough demo ready"),
        ("B6", "D2 \u00B7 08-10", "MVP sprint #2", "Polished prototype"),
        ("B7", "D2 \u00B7 10-12", "Dry-run + feedback", "Revised pitch deck"),
        ("B8", "D2 \u00B7 13-15", "Demo day + awards", "5-min pitch + award"),
    ]
    col_w = (CONTENT_W - Inches(0.45)) / 4
    row_h = Inches(1.45)
    for i, (code, time, title, artifact) in enumerate(blocks):
        r = i // 4
        c = i % 4
        x = MARGIN + (col_w + Inches(0.15)) * c
        y = Inches(1.55) + (row_h + Inches(0.15)) * r
        add_rect(s, x, y, col_w, row_h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, col_w, Inches(0.36), fill=ORANGE)
        add_text(s, x + Inches(0.15), y + Inches(0.04), Inches(0.5), Inches(0.3),
                 code, size=Pt(11), bold=True, color=WHITE)
        add_text(s, x + Inches(0.6), y + Inches(0.05), col_w - Inches(0.7), Inches(0.3),
                 time, size=Pt(9), color=WHITE, align=PP_ALIGN.RIGHT)
        add_text(s, x + Inches(0.15), y + Inches(0.45), col_w - Inches(0.25), Inches(0.5),
                 title, size=Pt(11), bold=True, color=BLACK, line_spacing=1.2)
        add_text(s, x + Inches(0.15), y + Inches(0.95), col_w - Inches(0.25), Inches(0.45),
                 "Out: " + artifact, size=Pt(9), italic=True, color=ORANGE, line_spacing=1.25)
    add_footer(s, "Part C \u00B7 Why 1N2D", 17)

    # Slide 18 — Coaching Cadence
    s = new_slide(prs)
    add_title_bar(s, "Three to four 1:1 sessions per team. Each session has a named purpose.",
                  eyebrow="Coaching cadence")
    sessions = [
        ("Session 1", "Block B4 \u00B7 ~15:30", "Problem + idea review",
         "Stress-test problem. Kill weak hypotheses. Assign riskiest assumption."),
        ("Session 2", "Block B5 \u00B7 ~20:30", "MVP direction",
         "Choose validation instrument. Scope to 90-min buildable."),
        ("Session 3", "Block B6 \u00B7 ~09:00", "Execution unblocking",
         "Fix blocker. Re-scope if off-track. Protect the demo."),
        ("Session 4", "Block B7 \u00B7 ~11:00", "Pitch coaching",
         "Deck order. Opening hook. Handle the hardest question."),
    ]
    y = Inches(1.55)
    row_h = Inches(0.78)
    for i, (label, when, title, body) in enumerate(sessions):
        ry = y + (row_h + Inches(0.08)) * i
        add_rect(s, MARGIN, ry, CONTENT_W, row_h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, MARGIN, ry, Inches(0.06), row_h, fill=ORANGE)
        add_text(s, MARGIN + Inches(0.2), ry + Inches(0.1), Inches(1.4), Inches(0.3),
                 label.upper(), size=Pt(10), bold=True, color=ORANGE)
        add_text(s, MARGIN + Inches(0.2), ry + Inches(0.36), Inches(1.4), Inches(0.3),
                 when, size=Pt(9), color=GRAY_BODY)
        add_text(s, MARGIN + Inches(1.7), ry + Inches(0.12), Inches(2.6), Inches(0.55),
                 title, size=Pt(12), bold=True, color=BLACK, line_spacing=1.2)
        add_text(s, MARGIN + Inches(4.4), ry + Inches(0.12), CONTENT_W - Inches(4.6), row_h - Inches(0.25),
                 body, size=Pt(10), color=GRAY_BODY, line_spacing=1.35)
    add_text(s, MARGIN, Inches(5.0), CONTENT_W, Inches(0.3),
             "Benchmark from SU Team: 800 total 1:1 sessions across 150 teams proves the density model.",
             size=Pt(10), italic=True, color=GRAY_BODY)
    add_footer(s, "Part C \u00B7 Why 1N2D", 18)

    # Slide 19 — 4-axis rubric
    s = new_slide(prs)
    add_title_bar(s, "Four axes, 25% each. Three judges. No hidden criteria.",
                  eyebrow="Assessment rubric (from su team)")
    axes = [
        ("25%", "Idea specificity", "Is the problem concrete? Is the insight non-obvious?"),
        ("25%", "BM completeness", "Who pays, how much, via what channel? Unit economics?"),
        ("25%", "Growth potential", "Can this 10x? What is the wedge + expansion path?"),
        ("25%", "Impact", "Market size \u00D7 social/economic ripple. Tied to Outcome."),
    ]
    col_w = (CONTENT_W - Inches(0.45)) / 4
    y = Inches(1.55)
    h = Inches(2.4)
    for i, (pct, title, desc) in enumerate(axes):
        x = MARGIN + (col_w + Inches(0.15)) * i
        add_rect(s, x, y, col_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.3), Inches(0.6),
                 pct, size=Pt(28), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.85), col_w - Inches(0.3), Inches(0.45),
                 title, size=Pt(12), bold=True, color=BLACK, line_spacing=1.2)
        add_text(s, x + Inches(0.2), y + Inches(1.35), col_w - Inches(0.3), h - Inches(1.45),
                 desc, size=Pt(9), color=GRAY_BODY, line_spacing=1.35)
    # Jury composition
    add_rect(s, MARGIN, Inches(4.2), CONTENT_W, Inches(0.7), fill=BLACK)
    add_rect(s, MARGIN, Inches(4.2), Inches(0.08), Inches(0.7), fill=ORANGE)
    add_text(s, MARGIN + Inches(0.25), Inches(4.3), CONTENT_W - Inches(0.4), Inches(0.25),
             "JURY COMPOSITION", size=Pt(9), bold=True, color=ORANGE)
    add_text(s, MARGIN + Inches(0.25), Inches(4.55), CONTENT_W - Inches(0.4), Inches(0.35),
             "3 people = 1 public-sector expert  +  1 private-sector expert  +  1 Underdogs coach.",
             size=Pt(11), color=WHITE)
    add_footer(s, "Part C \u00B7 Why 1N2D", 19)

    # =============================================================
    # PART D — Country Tracks (slides 20-36)
    # =============================================================

    # Slide 20 — Part D divider
    section_divider(prs,
        "Part D",
        "Same philosophy, different operating context.",
        "Five country tracks \u2014 Indonesia, Korea Inbound, Korea Outbound, Japan+Taiwan, India. Each gets three slides plus one comparison table.",
        20, "Part D \u00B7 Country Tracks")

    # ---- Indonesia slides 21-23 ----
    # 21: Overview
    s = new_slide(prs)
    add_title_bar(s, "Indonesia \u00B7 BSD City \u2014 a 500k-resident Living Lab for student founders.",
                  eyebrow="Track 01 \u00B7 Indonesia (BSD city)")
    cards = [
        {"tag": "Anchor", "title": "BSD City \u00B7 Tangerang",
         "body": ["10,000 ha planned city",
                  "500,000 residents",
                  "Indonesia's official Living Lab"]},
        {"tag": "Partners", "title": "Sinarmas + LLV + 7 universities",
         "body": ["Monash Indonesia \u00B7 Prasmul",
                  "BINUS \u00B7 Atma Jaya BSD",
                  "ITTS \u00B7 BSI \u00B7 UIC College BSD"]},
        {"tag": "Funnel", "title": "Stage 1 \u2192 2 \u2192 3",
         "body": ["S1: 1N2D hackathon (PoC)",
                  "S2: 3-4 month batch",
                  "S3: Alliance + joint fund"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.4))
    add_text(s, MARGIN, Inches(4.1), CONTENT_W, Inches(0.7),
             "Operating mode: Bilingual EN + Bahasa Indonesia. PoC discipline: \"Maximum Validation with Minimum Resources.\" Methodology: IMPACT + Local adaptation.",
             size=Pt(11), color=GRAY_BODY, line_spacing=1.4)
    add_footer(s, "Part D \u00B7 Indonesia", 21)

    # 22: 6 impact tracks
    s = new_slide(prs)
    add_title_bar(s, "Teams pick one of six BSD-anchored impact tracks.",
                  eyebrow="Indonesia \u00B7 problem tracks")
    tracks = [
        ("01", "Urban Mobility", "Traffic, last-mile, EV infrastructure"),
        ("02", "Health & Biomedical", "KEK ETKI special zone \u00B7 aging care"),
        ("03", "Green & Sustainability", "Green Habit since 2017 \u00B7 waste loops"),
        ("04", "Digital Inclusion", "Apple Academy \u00B7 Purwadhika \u00B7 access"),
        ("05", "UMKM & Local Economy", "65.5M MSMEs nationally \u00B7 B2B enablement"),
        ("06", "Education Access", "K-12 + HE \u00B7 language + cost barriers"),
    ]
    col_w = (CONTENT_W - Inches(0.3)) / 3
    row_h = Inches(1.3)
    for i, (num, title, desc) in enumerate(tracks):
        r = i // 3
        c = i % 3
        x = MARGIN + (col_w + Inches(0.15)) * c
        y = Inches(1.55) + (row_h + Inches(0.15)) * r
        add_rect(s, x, y, col_w, row_h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), row_h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.3), Inches(0.28),
                 num, size=Pt(10), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.42), col_w - Inches(0.3), Inches(0.4),
                 title, size=Pt(13), bold=True, color=BLACK, line_spacing=1.15)
        add_text(s, x + Inches(0.2), y + Inches(0.85), col_w - Inches(0.3), row_h - Inches(0.95),
                 desc, size=Pt(9), color=GRAY_BODY, line_spacing=1.3)
    add_footer(s, "Part D \u00B7 Indonesia", 22)

    # 23: Operating model
    s = new_slide(prs)
    add_title_bar(s, "Operating model: bilingual delivery, coach-partner boundary, field-lab posture.",
                  eyebrow="Indonesia \u00B7 operating model")
    cards = [
        {"tag": "Pre-event", "title": "2-3 weeks async",
         "body": ["Bilingual LMS (EN + ID)",
                  "Track-specific reading kit",
                  "Problem-selection exercise",
                  "Team-formation survey"]},
        {"tag": "On-site", "title": "1N2D at BSD venue",
         "body": ["Field visits in BSD",
                  "LLV / Sinarmas mentors",
                  "Customer interviews in-city",
                  "Bilingual dual-track coaching"]},
        {"tag": "Post-event", "title": "Funnel handoff",
         "body": ["Top teams \u2192 Stage 2 batch",
                  "University alumni network",
                  "LLV fund pipeline meeting",
                  "Report back to Sinarmas"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.5))
    add_text(s, MARGIN, Inches(4.2), CONTENT_W, Inches(0.6),
             "Coach-partner boundary rule: Underdogs coaches own methodology + 1:1 coaching; local partners own venue, judges, and follow-up fund intros. Write the boundary into the coach brief.",
             size=Pt(10), italic=True, color=GRAY_BODY, line_spacing=1.4)
    add_footer(s, "Part D \u00B7 Indonesia", 23)

    # ---- Korea Inbound slides 24-26 ----
    # 24: Overview
    s = new_slide(prs)
    add_title_bar(s, "Korea Inbound \u2014 foreign students starting up inside Korea.",
                  eyebrow="Track 02 \u00B7 Korea inbound")
    cards = [
        {"tag": "Target", "title": "International students in KR",
         "body": ["Undergrad + grad in KR universities",
                  "Intent: launch IN Korea",
                  "Often pre-team, pre-idea"]},
        {"tag": "Core challenge", "title": "K-market localization",
         "body": ["Language barrier",
                  "Regulatory + visa friction",
                  "Funding access (TIPS, KISED)"]},
        {"tag": "Theme", "title": "K-Market Entry for Foreign Founders",
         "body": ["Reverse glocal posture",
                  "Leverage alumni foreign-origin founders",
                  "Korean mentor pool essential"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.5))
    add_text(s, MARGIN, Inches(4.2), CONTENT_W, Inches(0.6),
             "Methodology: IMPACT + Small Business (for B2C/local) OR Global Entry reversed \u2014 they are the \"global\" founders entering KR.",
             size=Pt(11), color=GRAY_BODY, line_spacing=1.4)
    add_footer(s, "Part D \u00B7 Korea Inbound", 24)

    # 25: Pre-event + resources
    s = new_slide(prs)
    add_title_bar(s, "Pre-event teaches the ecosystem; on-site teaches the pitch.",
                  eyebrow="Korea inbound \u00B7 delivery")
    two_col_text(s, Inches(1.55), Inches(3.3),
                 "Pre-event LMS (2-3 weeks)",
                 ["Korean startup ecosystem basics",
                  "KOREA-VISA for founders",
                  "TIPS + KISED program overview",
                  "K-market consumer behavior",
                  "Alumni founder case studies"],
                 "Special resources on-site",
                 ["Korean mentor pool (local)",
                  "Foreign-origin alumni founders",
                  "Bilingual 1:1 coaching",
                  "K-market field interviews",
                  "VC intros (post-demo)"])
    add_footer(s, "Part D \u00B7 Korea Inbound", 25)

    # 26: Day flow
    s = new_slide(prs)
    add_title_bar(s, "Day flow emphasizes regulatory + cultural nuance every coach must surface.",
                  eyebrow="Korea inbound \u00B7 day flow")
    cards = [
        {"tag": "Day 1 AM", "title": "Problem in Korean context",
         "body": ["Interview 2 Koreans on block",
                  "Visa / biz-reg constraints in scope",
                  "Price / channel norms drilled"]},
        {"tag": "Day 1 PM", "title": "Prototype with K-UX",
         "body": ["Korean-language UI sketch",
                  "Payment method (Toss, KakaoPay)",
                  "Coach: localization gotchas"]},
        {"tag": "Day 2", "title": "Pitch for K-judges",
         "body": ["Korean + English dual deck",
                  "TIPS/KISED fit articulated",
                  "Judge: 1 KR VC + 1 policy + 1 Underdogs"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.7))
    add_footer(s, "Part D \u00B7 Korea Inbound", 26)

    # ---- Korea Outbound slides 27-29 ----
    # 27: Overview
    s = new_slide(prs)
    add_title_bar(s, "Korea Outbound \u2014 Korean teams born global from Day 1.",
                  eyebrow="Track 03 \u00B7 Korea outbound")
    cards = [
        {"tag": "Target", "title": "Korean teams expanding",
         "body": ["Already have KR traction",
                  "Eyeing JP / SE Asia / India / US",
                  "Need target-market discovery"]},
        {"tag": "Core challenge", "title": "Cross-border fit",
         "body": ["Target-market discovery",
                  "Local partner mapping",
                  "Cross-cultural team formation"]},
        {"tag": "Theme", "title": "Born Global from Day 1",
         "body": ["Methodology: Global Entry",
                  "Underdogs JP / India branches",
                  "520+ global partner network"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.5))
    add_footer(s, "Part D \u00B7 Korea Outbound", 27)

    # 28: Target market primers
    s = new_slide(prs)
    add_title_bar(s, "Pre-event primers by destination market. Teams pick their primer.",
                  eyebrow="Korea outbound \u00B7 market primers")
    markets = [
        ("JP", "Japan primer", "Regulatory caution \u00B7 trust-building \u00B7 sales cycles"),
        ("ID", "Indonesia primer", "500k-city living labs \u00B7 MSME \u00B7 Bahasa UX"),
        ("IN", "India primer", "Scale economics \u00B7 UPI / Aadhaar \u00B7 regional divide"),
        ("US", "US primer", "GTM velocity \u00B7 seed terms \u00B7 B2B SaaS patterns"),
    ]
    col_w = (CONTENT_W - Inches(0.45)) / 4
    y = Inches(1.55)
    h = Inches(2.2)
    for i, (code, title, desc) in enumerate(markets):
        x = MARGIN + (col_w + Inches(0.15)) * i
        add_rect(s, x, y, col_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.3), Inches(0.55),
                 code, size=Pt(28), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.8), col_w - Inches(0.3), Inches(0.5),
                 title, size=Pt(12), bold=True, color=BLACK, line_spacing=1.15)
        add_text(s, x + Inches(0.2), y + Inches(1.3), col_w - Inches(0.3), h - Inches(1.4),
                 desc, size=Pt(9), color=GRAY_BODY, line_spacing=1.35)
    add_text(s, MARGIN, Inches(4.1), CONTENT_W, Inches(0.6),
             "On-site: dedicated coach per target market. Demo day judge panel includes one investor from the team's target market whenever possible.",
             size=Pt(10), italic=True, color=GRAY_BODY, line_spacing=1.4)
    add_footer(s, "Part D \u00B7 Korea Outbound", 28)

    # 29: Partner leverage
    s = new_slide(prs)
    add_title_bar(s, "520+ global partners are the outbound team's unfair advantage.",
                  eyebrow="Korea outbound \u00B7 partner leverage")
    cards = [
        {"tag": "Before", "title": "Warm intros pre-event",
         "body": ["Partner list matched to target market",
                  "2-3 intro calls in pre-LMS weeks",
                  "Customer interview letters"]},
        {"tag": "During", "title": "Mentors on-site",
         "body": ["1 partner mentor per region",
                  "Judges drawn from network",
                  "Real-time market feedback"]},
        {"tag": "After", "title": "Landing support",
         "body": ["Underdogs JP / India landing",
                  "Local coach handoff",
                  "Batch program (Tier 3)"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.7))
    add_footer(s, "Part D \u00B7 Korea Outbound", 29)

    # ---- Japan + Taiwan slides 30-32 ----
    # 30: Overview
    s = new_slide(prs)
    add_title_bar(s, "Japan + Taiwan \u2014 East Asian social problems, one curated pool.",
                  eyebrow="Track 04 \u00B7 Japan + taiwan")
    cards = [
        {"tag": "Anchor JP", "title": "Musashino Univ.",
         "body": ["Faculty of Entrepreneurship",
                  "Tokyo 10x10x10 vision",
                  "Gov-private co-judged panels"]},
        {"tag": "Anchor TW", "title": "Taichung Social Innovation Lab",
         "body": ["7,000+ learner-contacts",
                  "350+ teams incubated",
                  "60+ award-winning teams"]},
        {"tag": "Frame", "title": "Common East-Asian problems",
         "body": ["Aging society",
                  "Environment + climate",
                  "Regional depopulation"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.5))
    add_text(s, MARGIN, Inches(4.2), CONTENT_W, Inches(0.5),
             "Methodology: Matching + Custom (social-enterprise oriented). Deliverable bar: BM with a validation plan attached.",
             size=Pt(11), color=GRAY_BODY, line_spacing=1.4)
    add_footer(s, "Part D \u00B7 Japan + Taiwan", 30)

    # 31: Problem pool
    s = new_slide(prs)
    add_title_bar(s, "Four problem categories. Teams pick one and commit.",
                  eyebrow="Japan + taiwan \u00B7 problem pool")
    pool = [
        ("01", "Aging Society", "Eldercare workforce \u00B7 loneliness \u00B7 dementia support"),
        ("02", "Environment", "Waste loops \u00B7 water \u00B7 urban biodiversity"),
        ("03", "Regional Depopulation", "Empty towns \u00B7 services \u00B7 young talent return"),
        ("04", "Climate Change", "Green GX \u00B7 disaster prep \u00B7 food resilience"),
    ]
    col_w = (CONTENT_W - Inches(0.45)) / 4
    y = Inches(1.55)
    h = Inches(2.5)
    for i, (num, title, desc) in enumerate(pool):
        x = MARGIN + (col_w + Inches(0.15)) * i
        add_rect(s, x, y, col_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.3), Inches(0.3),
                 num, size=Pt(10), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.45), col_w - Inches(0.3), Inches(0.8),
                 title, size=Pt(14), bold=True, color=BLACK, line_spacing=1.15)
        add_text(s, x + Inches(0.2), y + Inches(1.3), col_w - Inches(0.3), h - Inches(1.4),
                 desc, size=Pt(9), color=GRAY_BODY, line_spacing=1.4)
    add_footer(s, "Part D \u00B7 Japan + Taiwan", 31)

    # 32: SU Team pattern
    s = new_slide(prs)
    add_title_bar(s, "Inherit EASI LEAGUE's 4-axis rubric and 3-4 coach cadence.",
                  eyebrow="Japan + taiwan \u00B7 su team pattern")
    cards = [
        {"tag": "Judging", "title": "4-axis rubric",
         "body": ["Idea specificity",
                  "BM completeness",
                  "Growth potential",
                  "Impact (market + social)"]},
        {"tag": "Coaching", "title": "3-4 sessions / team",
         "body": ["Benchmark: 800 1:1 for 150 teams",
                  "Gov + private co-mentors",
                  "Local coach is table stakes"]},
        {"tag": "Bar", "title": "BM + validation plan",
         "body": ["No idea-only pitches",
                  "Customer-interview log",
                  "Next 30-day validation step"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.7))
    add_footer(s, "Part D \u00B7 Japan + Taiwan", 32)

    # ---- India slides 33-35 ----
    # 33: Overview
    s = new_slide(prs)
    add_title_bar(s, "India \u00B7 Action AI \u2014 launchpad, not course. Live URL by Day 1 end.",
                  eyebrow="Track 05 \u00B7 india (action ai)")
    cards = [
        {"tag": "Program", "title": "Action AI line",
         "body": ["Fin-Act track",
                  "AI Literacy track",
                  "Sold to universities"]},
        {"tag": "Philosophy", "title": "Many founders, not few unicorns",
         "body": ["One 5D per 1N2D",
                  "Chosen D: D4 Design",
                  "MVP prototyping focus"]},
        {"tag": "AI stack", "title": "Free-tier tools",
         "body": ["Gemini",
                  "Lovable",
                  "Claude Code"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(2.5))
    add_text(s, MARGIN, Inches(4.2), CONTENT_W, Inches(0.5),
             "Operating mode: \"Not a course, a launchpad.\" Every team ships a live product URL before Day 1 closes.",
             size=Pt(11), italic=True, color=GRAY_BODY, line_spacing=1.4)
    add_footer(s, "Part D \u00B7 India", 33)

    # 34: 6 tracks
    s = new_slide(prs)
    add_title_bar(s, "Six Action AI tracks. Teams pick one at registration.",
                  eyebrow="India \u00B7 action ai track pool")
    pool = [
        ("01", "Access & Regional Divide", "Language + connectivity + rural-urban gap"),
        ("02", "AI \u00D7 Senior Life", "Eldercare \u00B7 companionship \u00B7 health"),
        ("03", "Industrial Transformation", "Manufacturing \u00B7 logistics \u00B7 MSME tooling"),
        ("04", "Everyday Safety", "Women safety \u00B7 road \u00B7 emergency response"),
        ("05", "Climate & Circular", "Waste loops \u00B7 solar \u00B7 water resilience"),
        ("06", "Social Cohesion & Open", "Civic tech \u00B7 open data \u00B7 inclusion"),
    ]
    col_w = (CONTENT_W - Inches(0.3)) / 3
    row_h = Inches(1.3)
    for i, (num, title, desc) in enumerate(pool):
        r = i // 3
        c = i % 3
        x = MARGIN + (col_w + Inches(0.15)) * c
        y = Inches(1.55) + (row_h + Inches(0.15)) * r
        add_rect(s, x, y, col_w, row_h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), row_h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.3), Inches(0.25),
                 num, size=Pt(10), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.4), col_w - Inches(0.3), Inches(0.45),
                 title, size=Pt(12), bold=True, color=BLACK, line_spacing=1.15)
        add_text(s, x + Inches(0.2), y + Inches(0.88), col_w - Inches(0.3), row_h - Inches(0.98),
                 desc, size=Pt(9), color=GRAY_BODY, line_spacing=1.3)
    add_footer(s, "Part D \u00B7 India", 34)

    # 35: Funnel 20k to 10
    s = new_slide(prs)
    add_title_bar(s, "20,000 to 10: a performance funnel culminating in Korea Grand Showcase.",
                  eyebrow="India \u00B7 performance progression")
    stages = [
        ("20,000", "Registered"),
        ("1,000", "Qualified"),
        ("100", "Finalists"),
        ("10", "Grand Showcase"),
    ]
    seg_w = (CONTENT_W - Inches(0.45)) / 4
    y = Inches(1.7)
    h = Inches(2.1)
    for i, (num, label) in enumerate(stages):
        x = MARGIN + (seg_w + Inches(0.15)) * i
        add_rect(s, x, y, seg_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), h, fill=ORANGE)
        add_text(s, x + Inches(0.15), y + Inches(0.25), seg_w - Inches(0.25), Inches(0.9),
                 num, size=Pt(32), bold=True, color=ORANGE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.15), y + Inches(1.25), seg_w - Inches(0.25), Inches(0.7),
                 label, size=Pt(12), bold=True, color=BLACK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        if i < 3:
            ax = x + seg_w - Inches(0.03)
            add_text(s, ax, y + Inches(0.85), Inches(0.2), Inches(0.4),
                     "\u25B8", size=Pt(16), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(4.1), CONTENT_W, Inches(0.6),
             "Top 10 from India travel to Korea Grand Showcase. The 1N2D is the 1,000 \u2192 100 stage \u2014 build one MVP, ship live URL, pitch.",
             size=Pt(11), italic=True, color=GRAY_BODY, line_spacing=1.4)
    add_footer(s, "Part D \u00B7 India", 35)

    # 36: Comparison table
    s = new_slide(prs)
    add_title_bar(s, "Five tracks side by side \u2014 same skeleton, different dials.",
                  eyebrow="Comparison table")
    headers = ["Track", "Anchor", "Methodology", "Language", "Signature move"]
    rows = [
        ["Indonesia", "BSD City \u00B7 Sinarmas + 7 unis", "IMPACT + Local", "EN + Bahasa", "Living Lab field visits"],
        ["KR-Inbound", "Foreign students in KR", "IMPACT + Small Biz", "EN + Korean", "K-market interviews"],
        ["KR-Outbound", "520+ global partners", "Global Entry", "EN + KR + target", "Market primer pick"],
        ["JP + TW", "Musashino + Taichung", "Matching + Custom", "EN + JP + TW", "Gov-private co-judging"],
        ["India", "Action AI \u00B7 universities", "IMPACT + Custom", "EN + Hindi + reg.", "Live URL Day 1 end"],
    ]
    col_widths = [Inches(1.3), Inches(2.5), Inches(1.9), Inches(1.7), Inches(1.6)]
    y = Inches(1.5)
    # Header
    x = MARGIN
    add_rect(s, MARGIN, y, CONTENT_W, Inches(0.4), fill=BLACK)
    for i, h in enumerate(headers):
        add_text(s, x + Inches(0.12), y + Inches(0.08), col_widths[i], Inches(0.25),
                 h, size=Pt(10), bold=True, color=WHITE)
        x += col_widths[i]
    y += Inches(0.4)
    row_h = Inches(0.55)
    for ri, row in enumerate(rows):
        bg = CARD_BG if ri % 2 == 0 else WHITE
        add_rect(s, MARGIN, y, CONTENT_W, row_h, fill=bg, border=BORDER, border_w=0.3)
        x = MARGIN
        for ci, cell in enumerate(row):
            bold = ci == 0
            color = BLACK if ci == 0 else GRAY_BODY
            add_text(s, x + Inches(0.12), y + Inches(0.15), col_widths[ci] - Inches(0.15), row_h - Inches(0.2),
                     cell, size=Pt(9), bold=bold, color=color, line_spacing=1.2)
            x += col_widths[ci]
        y += row_h
    add_footer(s, "Part D \u00B7 Comparison", 36)

    # =============================================================
    # PART E — Facilitator Playbook (slides 37-42)
    # =============================================================

    # Slide 37 — Part E divider
    section_divider(prs,
        "Part E",
        "Facilitator playbook.",
        "Before, during, and after the 36 hours. Pre-event design, Day 1 + Day 2 timelines, coach cue cards, and post-event handoff.",
        37, "Part E \u00B7 Playbook")

    # Slide 38 — Pre-event LMS/VOD
    s = new_slide(prs)
    add_title_bar(s, "Pre-event is 2-3 weeks async. Concepts live here, not on-site.",
                  eyebrow="Pre-event \u00B7 lms / vod design")
    cards = [
        {"tag": "Week 1", "title": "Orient + problem",
         "body": ["2-3 videos (10-15 min)",
                  "Reading kit: 2-3 pieces",
                  "Exercise: problem one-pager",
                  "Deliverable: problem statement v0"]},
        {"tag": "Week 2", "title": "Method + tools",
         "body": ["3-4 videos on 5D + rubric",
                  "Reading: 1-2 track case studies",
                  "Exercise: tool walkthrough",
                  "Deliverable: tool-practice log"]},
        {"tag": "Week 3", "title": "Team + prep",
         "body": ["Team-formation survey",
                  "Interview 2 target users",
                  "Bring 1 draft idea + 1 backup",
                  "Deliverable: arrival packet"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(3.2))
    add_footer(s, "Part E \u00B7 Playbook", 38)

    # Slide 39 — Day 1 timeline
    s = new_slide(prs)
    add_title_bar(s, "Day 1 \u00B7 09:00 \u2192 22:00. Five blocks, each with a named artifact.",
                  eyebrow="Day 1 timeline")
    rows = [
        ("09:00-10:00", "B1 Opening + Icebreaker", "Team norms, intros, rubric briefing", "Signed team charter"),
        ("10:00-12:00", "B2 Problem deep-dive", "Track walk-through, interview review", "Problem statement v1"),
        ("13:00-15:00", "B3 Team + Ideation", "Team lock-in, idea branching", "Idea v1 + riskiest assumption"),
        ("15:00-18:00", "B4 MVP sprint #1 + 1:1 (Sess. 1)", "Build, test, coach", "First prototype"),
        ("19:00-22:00", "B5 Evening sprint + dinner + 1:1 (Sess. 2)", "Polish, rehearse, unblock", "Rough demo ready"),
    ]
    y = Inches(1.5)
    row_h = Inches(0.68)
    # Header
    add_rect(s, MARGIN, y, CONTENT_W, Inches(0.38), fill=BLACK)
    add_text(s, MARGIN + Inches(0.15), y + Inches(0.07), Inches(1.4), Inches(0.25),
             "Time", size=Pt(10), bold=True, color=WHITE)
    add_text(s, MARGIN + Inches(1.55), y + Inches(0.07), Inches(3.0), Inches(0.25),
             "Block", size=Pt(10), bold=True, color=WHITE)
    add_text(s, MARGIN + Inches(4.6), y + Inches(0.07), Inches(3.0), Inches(0.25),
             "What happens", size=Pt(10), bold=True, color=WHITE)
    add_text(s, MARGIN + Inches(7.65), y + Inches(0.07), Inches(1.35), Inches(0.25),
             "Artifact out", size=Pt(10), bold=True, color=WHITE)
    y += Inches(0.38)
    for i, (t, block, what, artifact) in enumerate(rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_rect(s, MARGIN, y, CONTENT_W, row_h, fill=bg, border=BORDER, border_w=0.3)
        add_text(s, MARGIN + Inches(0.15), y + Inches(0.15), Inches(1.4), row_h - Inches(0.2),
                 t, size=Pt(9), bold=True, color=ORANGE)
        add_text(s, MARGIN + Inches(1.55), y + Inches(0.15), Inches(3.0), row_h - Inches(0.2),
                 block, size=Pt(10), bold=True, color=BLACK, line_spacing=1.2)
        add_text(s, MARGIN + Inches(4.6), y + Inches(0.15), Inches(3.0), row_h - Inches(0.2),
                 what, size=Pt(9), color=GRAY_BODY, line_spacing=1.3)
        add_text(s, MARGIN + Inches(7.65), y + Inches(0.15), Inches(1.35), row_h - Inches(0.2),
                 artifact, size=Pt(9), italic=True, color=ORANGE, line_spacing=1.3)
        y += row_h
    add_footer(s, "Part E \u00B7 Playbook", 39)

    # Slide 40 — Day 2 timeline
    s = new_slide(prs)
    add_title_bar(s, "Day 2 \u00B7 08:00 \u2192 17:00. Polish, pitch, select, brief.",
                  eyebrow="Day 2 timeline")
    rows = [
        ("08:00-10:00", "B6 MVP sprint #2 + 1:1 (Sess. 3)", "Polish + demo prep", "Polished prototype"),
        ("10:00-12:00", "B7 Dry-run + coach feedback + 1:1 (Sess. 4)", "Rehearse, refine deck", "Revised pitch deck"),
        ("13:00-15:00", "B8 Demo Day (5-min pitch + Q&A)", "Public judging", "Scored pitch"),
        ("15:00-16:00", "Judging + awards ceremony", "Rubric tally + announcement", "Award list"),
        ("16:00-17:00", "Next-steps briefing", "Stage 2 batch onboarding", "Alumni handoff packet"),
    ]
    y = Inches(1.5)
    row_h = Inches(0.68)
    add_rect(s, MARGIN, y, CONTENT_W, Inches(0.38), fill=BLACK)
    add_text(s, MARGIN + Inches(0.15), y + Inches(0.07), Inches(1.4), Inches(0.25),
             "Time", size=Pt(10), bold=True, color=WHITE)
    add_text(s, MARGIN + Inches(1.55), y + Inches(0.07), Inches(3.3), Inches(0.25),
             "Block", size=Pt(10), bold=True, color=WHITE)
    add_text(s, MARGIN + Inches(4.9), y + Inches(0.07), Inches(2.7), Inches(0.25),
             "What happens", size=Pt(10), bold=True, color=WHITE)
    add_text(s, MARGIN + Inches(7.65), y + Inches(0.07), Inches(1.35), Inches(0.25),
             "Artifact out", size=Pt(10), bold=True, color=WHITE)
    y += Inches(0.38)
    for i, (t, block, what, artifact) in enumerate(rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        add_rect(s, MARGIN, y, CONTENT_W, row_h, fill=bg, border=BORDER, border_w=0.3)
        add_text(s, MARGIN + Inches(0.15), y + Inches(0.15), Inches(1.4), row_h - Inches(0.2),
                 t, size=Pt(9), bold=True, color=ORANGE)
        add_text(s, MARGIN + Inches(1.55), y + Inches(0.15), Inches(3.3), row_h - Inches(0.2),
                 block, size=Pt(10), bold=True, color=BLACK, line_spacing=1.2)
        add_text(s, MARGIN + Inches(4.9), y + Inches(0.15), Inches(2.7), row_h - Inches(0.2),
                 what, size=Pt(9), color=GRAY_BODY, line_spacing=1.3)
        add_text(s, MARGIN + Inches(7.65), y + Inches(0.15), Inches(1.35), row_h - Inches(0.2),
                 artifact, size=Pt(9), italic=True, color=ORANGE, line_spacing=1.3)
        y += row_h
    add_footer(s, "Part E \u00B7 Playbook", 40)

    # Slide 41 — Cue cards
    s = new_slide(prs)
    add_title_bar(s, "Coach cue cards \u2014 what to ask, what to avoid, per rubric axis.",
                  eyebrow="Coaching cue cards")
    cards = [
        {"tag": "Idea", "title": "Ask",
         "body": ["Who exactly has this problem?",
                  "What did you learn from 2 interviews?",
                  "What's the non-obvious wedge?",
                  "Avoid: \"cool idea, keep going\""]},
        {"tag": "BM", "title": "Ask",
         "body": ["Who pays, how much, how often?",
                  "What's the unit economics?",
                  "Name 3 distribution channels",
                  "Avoid: \"we'll figure revenue later\""]},
        {"tag": "Growth", "title": "Ask",
         "body": ["What's the 10x path?",
                  "Where does the wedge expand?",
                  "What breaks at 1,000 users?",
                  "Avoid: \"total addressable market is huge\""]},
        {"tag": "Impact", "title": "Ask",
         "body": ["Which Outcome layer moves?",
                  "Who benefits beyond the buyer?",
                  "How would you measure in 12 months?",
                  "Avoid: \"we help the world\""]},
    ]
    col_w = (CONTENT_W - Inches(0.45)) / 4
    y = Inches(1.55)
    h = Inches(3.2)
    for i, c in enumerate(cards):
        x = MARGIN + (col_w + Inches(0.15)) * i
        add_rect(s, x, y, col_w, h, fill=CARD_BG, border=BORDER, border_w=0.5)
        add_rect(s, x, y, Inches(0.06), h, fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.18), col_w - Inches(0.3), Inches(0.22),
                 c["tag"].upper(), size=Pt(9), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.45), col_w - Inches(0.3), Inches(0.35),
                 c["title"], size=Pt(13), bold=True, color=BLACK)
        add_bullets(s, x + Inches(0.2), y + Inches(0.9), col_w - Inches(0.3), h - Inches(1.0),
                    c["body"], size=Pt(9), color=GRAY_BODY, spacing=1.3)
    add_footer(s, "Part E \u00B7 Playbook", 41)

    # Slide 42 — Post-event
    s = new_slide(prs)
    add_title_bar(s, "Post-event is where Stage-2 pipeline is won or lost.",
                  eyebrow="Post-event \u00B7 demo day \u2192 batch handoff")
    cards = [
        {"tag": "Week 0", "title": "Close out",
         "body": ["1-pager per team",
                  "Scorecard published",
                  "Photo + highlights reel",
                  "Debrief: coaches + partners"]},
        {"tag": "Weeks 2-4", "title": "Batch gate",
         "body": ["Top 10-20% selected",
                  "Stage 2 onboarding calls",
                  "Contract + coach assignment",
                  "First milestone scoped"]},
        {"tag": "Weeks 4-8", "title": "Alumni network",
         "body": ["All teams \u2192 alumni list",
                  "Quarterly check-in",
                  "Alumni-to-alumni intros",
                  "Next-cohort recruiters"]},
    ]
    three_column_cards(s, Inches(1.55), cards, card_h=Inches(3.2))
    add_footer(s, "Part E \u00B7 Playbook", 42)

    # =============================================================
    # PART F — Close (slides 43-44)
    # =============================================================

    # Slide 43 — Key Takeaways
    s = new_slide(prs)
    add_title_bar(s, "Five takeaways to carry into your next hackathon design.",
                  eyebrow="Key takeaways")
    takeaways = [
        ("01", "Philosophy first", "The 6 principles + 5D + Impact Value Chain are non-negotiable. Methodology picks itself after."),
        ("02", "Funnel is the format", "1N2D alone is not the program. Pre-LMS (2-3 wks) + 1N2D + follow-up batch is."),
        ("03", "One D per hackathon", "Pick a single 5D competency to install. Attempting all five in 36 hours fails."),
        ("04", "Coaching density wins", "3-4 1:1s per team, each with a named purpose, drives the quality delta."),
        ("05", "Localize the dials, keep the skeleton", "5 country tracks share one 8-block template. Only the content layer changes."),
    ]
    y = Inches(1.5)
    row_h = Inches(0.68)
    for i, (num, title, desc) in enumerate(takeaways):
        ry = y + (row_h + Inches(0.05)) * i
        add_rect(s, MARGIN, ry, CONTENT_W, row_h, fill=CARD_BG, border=BORDER, border_w=0.3)
        add_rect(s, MARGIN, ry, Inches(0.06), row_h, fill=ORANGE)
        add_text(s, MARGIN + Inches(0.22), ry + Inches(0.18), Inches(0.5), Inches(0.35),
                 num, size=Pt(14), bold=True, color=ORANGE)
        add_text(s, MARGIN + Inches(0.85), ry + Inches(0.1), Inches(2.6), Inches(0.5),
                 title, size=Pt(12), bold=True, color=BLACK, line_spacing=1.2)
        add_text(s, MARGIN + Inches(3.55), ry + Inches(0.12), CONTENT_W - Inches(3.75), row_h - Inches(0.2),
                 desc, size=Pt(10), color=GRAY_BODY, line_spacing=1.35)
    add_text(s, MARGIN, Inches(5.0), CONTENT_W, Inches(0.3),
             "Assignment: produce one country-track run-sheet using the provided template. See deliverables/assignments/.",
             size=Pt(10), italic=True, color=ORANGE, bold=True)
    add_footer(s, "Part F \u00B7 Close", 43)

    # Slide 44 — Thank you / contact
    s = new_slide(prs)
    # Dark closing slide
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=BLACK)
    add_rect(s, Inches(0.5), Inches(1.2), Inches(0.08), Inches(2.5), fill=ORANGE)
    add_text(s, Inches(0.85), Inches(1.1), Inches(8.3), Inches(0.4),
             "THANK YOU", size=Pt(11), bold=True, color=ORANGE)
    add_text(s, Inches(0.85), Inches(1.45), Inches(8.3), Inches(1.5),
             "Go design\nthe next 36 hours.", size=Pt(40), bold=True, color=WHITE, line_spacing=1.05)
    # Contact / resource grid
    contacts = [
        ("Contact", "udpb@udimpact.ai"),
        ("Website", "underdogs.global"),
        ("Coach finder", "coach.udimpact.ai"),
        ("Planning agent", "app.udimpact.ai"),
    ]
    col_w = (CONTENT_W - Inches(0.45)) / 4
    y = Inches(3.8)
    for i, (label, value) in enumerate(contacts):
        x = MARGIN + (col_w + Inches(0.15)) * i
        add_rect(s, x, y, col_w, Inches(0.85), fill=DARK)
        add_rect(s, x, y, Inches(0.06), Inches(0.85), fill=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.1), col_w - Inches(0.3), Inches(0.25),
                 label.upper(), size=Pt(9), bold=True, color=ORANGE)
        add_text(s, x + Inches(0.2), y + Inches(0.38), col_w - Inches(0.3), Inches(0.4),
                 value, size=Pt(11), bold=True, color=WHITE)
    # Footer page + copyright
    add_text(s, Inches(8.5), Inches(5.15), Inches(1.0), Inches(0.2),
             "44", size=TINY_SIZE, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(0.5), Inches(5.15), Inches(5.0), Inches(0.2),
             "Part F \u00B7 Close", size=TINY_SIZE, color=RGBColor(0x99, 0x99, 0x99))
    add_text(s, Inches(0.5), Inches(5.32), Inches(9.0), Inches(0.2),
             COPYRIGHT, size=MICRO_SIZE, color=RGBColor(0x99, 0x99, 0x99),
             align=PP_ALIGN.RIGHT)

    # ---------- save ----------
    out = Path(r"C:\Users\USER\projects\ud-ops-workspace\.claude\worktrees\amazing-khorana-50ddb7\deliverables\lecture-deck.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved {out}")
    print(f"Slide count: {len(prs.slides)}")
    print(f"Size: {out.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    build()
