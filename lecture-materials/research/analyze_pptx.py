"""Analyze AX Consultant Orientation PPTX for layout blueprint extraction."""
import sys
from collections import Counter, defaultdict
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX_PATH = r"C:\Users\USER\Desktop\AX컨설턴트_Orientation_260414.pptx"

def emu_to_inches(v):
    return round(v / 914400, 2) if v is not None else None

def emu_to_px(v):
    # 96 DPI
    return round(v / 914400 * 96) if v is not None else None

def rgb_hex(color):
    try:
        if color is None:
            return None
        if hasattr(color, 'rgb') and color.rgb is not None:
            return str(color.rgb)
    except Exception:
        return None
    return None

def extract_color_from_fill(fill):
    try:
        if fill.type == 1:  # solid
            return rgb_hex(fill.fore_color)
    except Exception:
        pass
    return None

def extract_font_info(run):
    font = run.font
    name = font.name
    try:
        size = font.size.pt if font.size else None
    except Exception:
        size = None
    color = None
    try:
        if font.color and font.color.type is not None:
            color = rgb_hex(font.color)
    except Exception:
        pass
    bold = font.bold
    return name, size, color, bold

def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from walk_shapes(shape.shapes)
            except Exception:
                pass

def main():
    prs = Presentation(PPTX_PATH)
    out = []
    out.append("# AX Consultant Orientation — Layout Blueprint Analysis\n")
    out.append(f"**Source:** `{PPTX_PATH}`\n")

    # Metadata
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    out.append("\n## 1. Deck Metadata\n")
    out.append(f"- **Total slides:** {len(prs.slides)}")
    out.append(f"- **Slide size (EMU):** {slide_w} x {slide_h}")
    out.append(f"- **Slide size (inches):** {emu_to_inches(slide_w)} x {emu_to_inches(slide_h)}")
    out.append(f"- **Slide size (px @ 96dpi):** {emu_to_px(slide_w)} x {emu_to_px(slide_h)}")
    aspect = slide_w / slide_h
    ratio = "16:9" if abs(aspect - 16/9) < 0.05 else ("4:3" if abs(aspect - 4/3) < 0.05 else f"{aspect:.3f}:1")
    out.append(f"- **Aspect ratio:** {ratio}")

    # Core properties
    try:
        cp = prs.core_properties
        out.append(f"- **Title:** {cp.title}")
        out.append(f"- **Author:** {cp.author}")
        out.append(f"- **Created:** {cp.created}")
        out.append(f"- **Modified:** {cp.modified}")
    except Exception as e:
        out.append(f"- Core properties read error: {e}")

    # Theme name (peek into xml)
    try:
        theme_part = prs.slide_masters[0].element
        # Try to get theme via related parts
        theme = None
        for rel in prs.slide_masters[0].part.rels.values():
            if "theme" in rel.reltype:
                theme_xml = rel.target_part.blob.decode("utf-8", errors="ignore")
                import re
                m = re.search(r'<a:theme[^>]*name="([^"]+)"', theme_xml)
                if m:
                    theme = m.group(1)
                break
        out.append(f"- **Theme name:** {theme or '(not found)'}")
    except Exception as e:
        out.append(f"- Theme read error: {e}")

    # Layouts
    out.append("\n## 2. Slide Masters & Layouts Defined\n")
    for master_i, master in enumerate(prs.slide_masters):
        out.append(f"### Master {master_i+1}")
        for layout in master.slide_layouts:
            out.append(f"  - Layout: **{layout.name}**")

    # Walk all slides
    colors = Counter()
    font_names = Counter()
    title_sizes = Counter()
    body_sizes = Counter()
    all_sizes = Counter()
    layout_usage = Counter()
    slide_reports = []
    slide_titles = []

    for i, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name
        layout_usage[layout_name] += 1

        shapes_info = []
        title_text = None
        first_text_size = None

        for shape in walk_shapes(slide.shapes):
            info = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left_in": emu_to_inches(shape.left),
                "top_in": emu_to_inches(shape.top),
                "width_in": emu_to_inches(shape.width),
                "height_in": emu_to_inches(shape.height),
                "has_text": False,
                "text_sample": None,
                "is_title": False,
                "is_placeholder": False,
                "ph_type": None,
                "fill_color": None,
                "font_name": None,
                "font_size": None,
                "font_color": None,
                "is_picture": shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
            }

            # Placeholder info
            if shape.is_placeholder:
                info["is_placeholder"] = True
                try:
                    info["ph_type"] = str(shape.placeholder_format.type)
                    if shape.placeholder_format.idx == 0 or "TITLE" in info["ph_type"]:
                        info["is_title"] = True
                except Exception:
                    pass

            # Fill color
            try:
                if hasattr(shape, 'fill'):
                    c = extract_color_from_fill(shape.fill)
                    if c:
                        info["fill_color"] = c
                        colors[c] += 1
            except Exception:
                pass

            # Text
            if shape.has_text_frame:
                tf = shape.text_frame
                txt_parts = []
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            txt_parts.append(run.text)
                        name, size, color, bold = extract_font_info(run)
                        if name:
                            font_names[name] += 1
                        if size:
                            all_sizes[size] += 1
                            if info["is_title"]:
                                title_sizes[size] += 1
                            else:
                                body_sizes[size] += 1
                            if info["font_size"] is None:
                                info["font_size"] = size
                        if name and info["font_name"] is None:
                            info["font_name"] = name
                        if color:
                            colors[color] += 1
                            if info["font_color"] is None:
                                info["font_color"] = color
                full_text = " ".join(txt_parts).strip()
                if full_text:
                    info["has_text"] = True
                    info["text_sample"] = full_text[:80]
                    if info["is_title"] and title_text is None:
                        title_text = full_text

            shapes_info.append(info)

        if title_text is None:
            # Fallback: use first large text
            for s in shapes_info:
                if s["has_text"] and s["font_size"] and s["font_size"] >= 24:
                    title_text = s["text_sample"]
                    break

        slide_titles.append(title_text or f"(slide {i} — no clear title)")
        slide_reports.append({
            "index": i,
            "layout": layout_name,
            "title": title_text,
            "shape_count": len(shapes_info),
            "shapes": shapes_info,
        })

    # Palette
    out.append("\n## 3. Color Palette\n")
    top_colors = colors.most_common(15)
    out.append("Top colors by frequency (fills + text):\n")
    for c, n in top_colors:
        out.append(f"- `#{c}` — {n} occurrences")

    # Fonts
    out.append("\n## 4. Fonts & Typography\n")
    out.append("**Font families used:**")
    for f, n in font_names.most_common():
        out.append(f"- `{f}` — {n} runs")
    out.append("\n**Font sizes (all runs):**")
    for sz, n in sorted(all_sizes.items(), key=lambda x: -x[1])[:15]:
        out.append(f"- {sz} pt — {n} runs")
    out.append("\n**Title-like sizes (placeholder=TITLE):**")
    for sz, n in sorted(title_sizes.items(), key=lambda x: -x[1])[:10]:
        out.append(f"- {sz} pt — {n} runs")
    out.append("\n**Body-like sizes:**")
    for sz, n in sorted(body_sizes.items(), key=lambda x: -x[1])[:10]:
        out.append(f"- {sz} pt — {n} runs")

    # Layout usage
    out.append("\n## 5. Layout Usage\n")
    for layout, n in layout_usage.most_common():
        out.append(f"- **{layout}** — used in {n} slides")

    # Per-slide breakdown
    out.append("\n## 6. Per-Slide Breakdown\n")
    for r in slide_reports:
        out.append(f"\n### Slide {r['index']} — Layout: `{r['layout']}`")
        out.append(f"- **Title:** {r['title']}")
        out.append(f"- **Shape count:** {r['shape_count']}")
        # Only show key shapes (text + picture)
        key = [s for s in r["shapes"] if s["has_text"] or s["is_picture"]]
        for s in key[:12]:
            pos = f"[L={s['left_in']}\" T={s['top_in']}\" W={s['width_in']}\" H={s['height_in']}\"]"
            if s["is_picture"]:
                out.append(f"  - PICTURE {s['name']} {pos}")
            else:
                font = f"{s['font_name'] or '?'} {s['font_size'] or '?'}pt"
                color = f" color=#{s['font_color']}" if s['font_color'] else ""
                fill = f" fill=#{s['fill_color']}" if s['fill_color'] else ""
                ph = f" [{s['ph_type']}]" if s['is_placeholder'] else ""
                sample = s["text_sample"] or ""
                out.append(f"  - TEXT{ph} {pos} {font}{color}{fill} — \"{sample}\"")

    # Titles TOC
    out.append("\n## 7. Table of Contents (by slide titles)\n")
    for i, t in enumerate(slide_titles, start=1):
        out.append(f"{i}. {t}")

    report = "\n".join(out)
    with open(r"C:\Users\USER\projects\ud-ops-workspace\.claude\worktrees\amazing-khorana-50ddb7\scratch\raw-analysis.md", "w", encoding="utf-8") as f:
        f.write(report)
    print("WROTE raw-analysis.md")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {emu_to_inches(slide_w)} x {emu_to_inches(slide_h)} inches")
    print(f"Layouts used: {dict(layout_usage)}")
    print(f"Top fonts: {font_names.most_common(5)}")
    print(f"Top colors: {colors.most_common(10)}")

if __name__ == "__main__":
    main()
